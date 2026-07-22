"""
商談用資料 / 店舗POP の PowerPoint (pptx) 生成（issue 202607221128）。

選択中の1店舗について、地図画像（``lib.static_map.render_static_map`` の出力）を
テンプレートの画像プレースホルダーに貼り付けた pptx を1枚生成する。

テンプレートは Databricks の Unity Catalog Volume から取得する（``lib/volume.py`` と
同じ WorkspaceClient パターン）。ローカル開発など Volume にアクセスできない場合は
リポジトリ同梱の ``images/template.pptx`` にフォールバックする。

Functions
---------
load_template_bytes(kind)               -- Volume からテンプレ取得（失敗時ローカルfallback）
load_caption_formats()                  -- [pptx] の定型文3種（file <- 上書き <- 既定）を返す
load_caption(store)                     -- config の定型文に小売店名称を差し込んだ文字列を返す
load_dated_caption(fmt_key, ts)         -- データ更新日時 ts を定型文へ差し込んだ文字列を返す
save_caption_formats(values)            -- 定型文を databricks_config.toml へ保存（読取専用時は例外）
build_store_pptx(template_bytes, png, caption_texts)
                                        -- 画像プレースホルダーに地図PNGを挿入し、テキスト
                                           プレースホルダーへ各 caption を入れて pptx bytes を返す
"""

import io
import logging
import re
import tomllib

logger = logging.getLogger(__name__)

_CONFIG_PATH = "config/databricks_config.toml"

# kind -> config キー（テンプレートファイル名）
_TEMPLATE_KEY = {
    "shoudan": "shoudan_template",
    "pop": "pop_template",
}

# テキストプレースホルダーへ入れる定型文（キー -> 既定文）。テンプレのテキストプレースホルダーを
# idx 昇順に並べ、この順で対応付ける（issue 202607221450 / 202607221705）。
# {store} は小売店名称、{year}/{month}/{day} はデータ最終更新日（JST）に置換する。
_CAPTION_DEFAULTS: dict[str, str] = {
    "store_caption_format": "{store}",
    "map_status_caption_format": "※地図、店舗状況は{year}年{month}月{day}日時点",
    "activity_caption_format": "※地図中の園は{year}年に啓発活動を実施いただいた園となります",
}

# 読取専用FS（Databricks Apps 等）でのプロセス内オーバーライド（lib/colors.py と同型）。
# save 失敗時に apply_caption_overrides() で差し込み、その場の生成に即時反映させる。
_caption_overrides: dict[str, str] = {}


def _load_pptx_config() -> dict:
    """Load the [pptx] section from config/databricks_config.toml."""
    try:
        with open(_CONFIG_PATH, "rb") as f:
            data = tomllib.load(f)
        return data.get("pptx", {})
    except FileNotFoundError:
        return {}


def load_template_bytes(kind: str) -> bytes:
    """Return the template pptx bytes for *kind* ("shoudan" or "pop").

    まず ``template_dir/{kind のテンプレファイル名}`` を Volume から取得し、失敗した場合は
    ``local_fallback``（既定 ``images/template.pptx``）を読む。Databricks 上では
    WorkspaceClient が自動でクレデンシャルを取得する。
    """
    if kind not in _TEMPLATE_KEY:
        raise ValueError(f"unknown template kind: {kind!r}")

    config = _load_pptx_config()
    template_dir = config.get("template_dir", "")
    filename = config.get(_TEMPLATE_KEY[kind], "")
    local_fallback = config.get("local_fallback", "images/template.pptx")

    if template_dir and filename:
        path = f"{template_dir.rstrip('/')}/{filename}"
        try:
            from databricks.sdk import WorkspaceClient  # noqa: PLC0415

            w = WorkspaceClient()
            resp = w.files.download(path)
            data = resp.contents.read()
            logger.info("テンプレートを取得しました: %s", path)
            return data
        except Exception as e:  # noqa: BLE001
            logger.warning(
                "Volume からのテンプレ取得に失敗（%s）。ローカルfallbackを使用: %s",
                path,
                e,
            )

    with open(local_fallback, "rb") as f:
        return f.read()


def default_caption_formats() -> dict[str, str]:
    """組み込み既定の定型文3種の複製を返す（テーマ設定ページの「既定に戻す」用）。"""
    return dict(_CAPTION_DEFAULTS)


def load_caption_formats() -> dict[str, str]:
    """テキストプレースホルダー定型文3種を返す（既定 <- config ファイル <- 上書き）。

    値の出所は ``config/databricks_config.toml`` の ``[pptx]`` と、読取専用FS 用の
    プロセス内オーバーライド。未設定・空文字のキーは既定（``_CAPTION_DEFAULTS``）を使う。
    """
    config = _load_pptx_config()
    formats: dict[str, str] = {}
    for key, default in _CAPTION_DEFAULTS.items():
        value = _caption_overrides.get(key) or config.get(key) or default
        formats[key] = value
    return formats


def apply_caption_overrides(values: dict) -> None:
    """定型文のプロセス内オーバーライドを差し込む（読取専用FS での即時反映用）。"""
    for key in _CAPTION_DEFAULTS:
        if key in values and isinstance(values[key], str):
            _caption_overrides[key] = values[key]


def load_caption(store: str | None) -> str:
    """選択中の小売店名称 *store* を定型文に差し込んだキャプション文字列を返す。

    定型文は ``[pptx] store_caption_format``（既定 ``"{store}"``）。``{store}`` が
    小売店名称に置換される。*store* が空/None のときは空文字を返す（テキスト挿入なし）。
    """
    if not store:
        return ""
    fmt = load_caption_formats()["store_caption_format"]
    try:
        return fmt.format(store=store)
    except (KeyError, IndexError) as e:  # noqa: PERF203
        logger.warning("store_caption_format の書式が不正です（%r）: %s", fmt, e)
        return store


def load_dated_caption(fmt_key: str, ts) -> str:
    """データ更新日時 *ts* を ``[pptx] {fmt_key}`` の定型文へ差し込んで返す。

    *ts* は年月日を持つオブジェクト（``pd.Timestamp``/``datetime``）または None。
    ``{year}/{month}/{day}`` を JST の年月日（ゼロ埋めしない）に置換する。*ts* が None
    のときは空文字を返す（日時未取得ならキャプションを挿入しない）。書式不正時も空文字。
    """
    if ts is None:
        return ""
    fmt = load_caption_formats().get(fmt_key, "")
    if not fmt:
        return ""
    try:
        return fmt.format(year=ts.year, month=ts.month, day=ts.day)
    except (KeyError, IndexError, ValueError) as e:  # noqa: PERF203
        logger.warning("%s の書式が不正です（%r）: %s", fmt_key, fmt, e)
        return ""


def _escape_toml_basic(value: str) -> str:
    """TOML 基本文字列用に ``\\`` と ``"`` をエスケープする。"""
    return value.replace("\\", "\\\\").replace('"', '\\"')


def patched_config_text(values: dict) -> str:
    """現 ``databricks_config.toml`` の定型文キーを *values* で置換した全文を返す。

    ``[pptx]`` セクション以外は保持し、対象キー行のみ差し替える（無ければ ``[pptx]`` 直下へ
    追記）。他セクションやコメントを壊さずに保存・ダウンロードするための最小パッチ。
    """
    try:
        with open(_CONFIG_PATH, encoding="utf-8") as f:
            text = f.read()
    except FileNotFoundError:
        text = "[pptx]\n"

    for key in _CAPTION_DEFAULTS:
        if key not in values or not isinstance(values[key], str):
            continue
        line = f'{key} = "{_escape_toml_basic(values[key])}"'
        pattern = re.compile(rf"^{re.escape(key)}\s*=.*$", re.MULTILINE)
        if pattern.search(text):
            text = pattern.sub(line, text, count=1)
        elif re.search(r"^\[pptx\]", text, re.MULTILINE):
            text = re.sub(
                r"(^\[pptx\][^\n]*\n)", rf"\1{line}\n", text, count=1, flags=re.MULTILINE
            )
        else:
            text = text.rstrip("\n") + f"\n\n[pptx]\n{line}\n"
    return text


def save_caption_formats(values: dict) -> str:
    """定型文 *values* を ``databricks_config.toml`` へ保存し、パスを返す。

    書き込み失敗（読取専用FS 等）時は例外を送出する（呼び出し側で処理）。
    """
    text = patched_config_text(values)
    with open(_CONFIG_PATH, "w", encoding="utf-8") as f:
        f.write(text)
    return _CONFIG_PATH


def _target_box(slide, prs):
    """画像を収める矩形 (left, top, width, height) を返し、対象プレースホルダーを削除する。

    画像プレースホルダー（``insert_picture`` 可）を最優先し、無ければ最初のプレースホルダー、
    それも無ければスライド全面を対象にする。プレースホルダーは空枠が残らないよう削除する。
    """
    picture_ph = None
    for ph in slide.placeholders:
        if hasattr(ph, "insert_picture"):
            picture_ph = ph
            break
    target = picture_ph or (slide.placeholders[0] if len(slide.placeholders) else None)
    if target is not None:
        box = (target.left, target.top, target.width, target.height)
        target._element.getparent().remove(target._element)
        return box
    return (0, 0, prs.slide_width, prs.slide_height)


def _fill_text_placeholders(slide, texts: list[str]) -> None:
    """スライドのテキストプレースホルダーへ *texts* を idx 昇順で書き込む。

    画像プレースホルダー（``insert_picture`` 可）は対象外とし、``text_frame`` を持つ
    プレースホルダーを idx 昇順（安定順）に並べ、*texts* の先頭から順に対応付ける。idx を
    ハードコードせず並び順で割り当てるため、テンプレによって idx 値が異なっても機能する。
    空文字の要素は該当枠を空のまま残す。枠数より *texts* が多い分は warning でスキップする。
    """
    text_phs = [
        ph
        for ph in slide.placeholders
        if not hasattr(ph, "insert_picture") and ph.has_text_frame
    ]
    text_phs.sort(key=lambda ph: ph.placeholder_format.idx)

    for i, text in enumerate(texts):
        if not text:
            continue
        if i >= len(text_phs):
            logger.warning(
                "テキストプレースホルダーが不足しています（必要 %d, テンプレ %d）: caption=%r",
                len(texts), len(text_phs), text,
            )
            break
        text_phs[i].text_frame.text = text


def build_store_pptx(
    template_bytes: bytes,
    map_png_bytes: bytes,
    caption_texts: list[str] | None = None,
) -> bytes:
    """Place *map_png_bytes* into the template's picture placeholder → pptx bytes.

    地図PNGを 1 枚目スライドの画像プレースホルダー位置へ貼り付ける。プレースホルダーの
    ``insert_picture`` は画像をプレースホルダーのアスペクト比にクロップしてしまい、対話地図と
    体裁が合わなくなる（issue 202607221245）。そこで、プレースホルダーの矩形内に **アスペクト比を
    保ったまま**収まるよう ``add_picture`` で中央配置する（元プレースホルダーは削除）。

    *caption_texts* が指定された場合は、テキストプレースホルダーへ idx 昇順に各文字列を
    書き込む（issue 202607221450 / 202607221705）。小売店名称・地図/店舗状況の時点・啓発活動年
    の3枠を想定。画像プレースホルダーとは別枠のため画像貼り付けと干渉しない。
    """
    from PIL import Image  # noqa: PLC0415
    from pptx import Presentation  # noqa: PLC0415

    prs = Presentation(io.BytesIO(template_bytes))
    slide = prs.slides[0]

    if caption_texts:
        _fill_text_placeholders(slide, caption_texts)

    box_left, box_top, box_w, box_h = _target_box(slide, prs)

    # 画像のアスペクト比を保って矩形にフィット（レターボックス）し、矩形中央へ配置する。
    with Image.open(io.BytesIO(map_png_bytes)) as im:
        img_w, img_h = im.size
    scale = min(box_w / img_w, box_h / img_h)
    draw_w = int(round(img_w * scale))
    draw_h = int(round(img_h * scale))
    left = int(box_left + (box_w - draw_w) / 2)
    top = int(box_top + (box_h - draw_h) / 2)

    slide.shapes.add_picture(
        io.BytesIO(map_png_bytes), left, top, width=draw_w, height=draw_h
    )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
